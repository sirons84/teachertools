# -*- coding: utf-8 -*-
"""
AI 토론 수업 설계 - PT용 PPT 생성 스크립트.
실행: python scripts/generate-pitch-pptx.py
출력: docs/AI-토론-수업-설계.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── 색상 팔레트 (앱 테마) ───────────────────────────────────────────
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_LIGHT = RGBColor(0xE0, 0xE7, 0xFF)
INDIGO_DARK = RGBColor(0x36, 0x30, 0xA3)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
BG = RGBColor(0xF8, 0xFA, 0xFC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

KOR_FONT = "맑은 고딕"


def set_run(run, text, font_size=18, bold=False, color=DARK, font=KOR_FONT):
    run.text = text
    run.font.name = font
    # set east-asian font too
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:eastAsianFont")) if False else None  # placeholder; python-pptx doesn't expose directly
    # We'll set ea via direct xml:
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is not None:
            rPr.remove(el)
        child = etree.SubElement(rPr, qn(f"a:{tag}"))
        child.set("typeface", font)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, font_size=font_size, bold=bold, color=color)
    return tb


def add_bullets(slide, left, top, width, height, bullets, font_size=18, color=DARK, line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(item, tuple):
            bullet_char, body = item
            set_run(run, f"{bullet_char}  {body}", font_size=font_size, color=color)
        else:
            set_run(run, f"•  {item}", font_size=font_size, color=color)
    return tb


def add_rect(slide, left, top, width, height, fill=WHITE, line_color=LIGHT_GRAY, line_width=1.0, rounded=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width)
    if not shadow:
        # remove default shadow
        spPr = shp.fill._xPr  # noqa
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    shp.shadow.inherit = False
    return shp


def put_text_in_shape(shape, text, font_size=14, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(40000)
    tf.margin_right = Emu(40000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    set_run(run, text, font_size=font_size, bold=bold, color=color)


def add_label_rect(slide, left, top, width, height, text, fill=INDIGO, text_color=WHITE, font_size=13, bold=True, rounded=True):
    shp = add_rect(slide, left, top, width, height, fill=fill, line_color=fill, rounded=rounded)
    put_text_in_shape(shp, text, font_size=font_size, bold=bold, color=text_color)
    return shp


def add_connector(slide, x1, y1, x2, y2, color=GRAY, width=1.5, arrow=True):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        lnXml = line.line._get_or_add_ln()
        tail = etree.SubElement(lnXml, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("h", "med")
    return line


def add_header_bar(slide, title, subtitle=None):
    """상단 얇은 인디고 바 + 제목."""
    # 상단 얇은 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INDIGO
    bar.line.fill.background()

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                title, font_size=32, bold=True, color=DARK)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
                    subtitle, font_size=16, color=GRAY)


def add_page_footer(slide, page_num, total):
    add_textbox(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
                f"{page_num} / {total}", font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.6), Inches(7.1), Inches(6), Inches(0.3),
                "TeacherTools · AI 토론 수업 설계", font_size=10, color=GRAY)


# ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

TOTAL = 12

# ── Slide 1 : Cover ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
# 배경 그라데이션 느낌 (좌측 얇은 큰 색 바)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.6), Inches(7.5))
accent.fill.solid()
accent.fill.fore_color.rgb = INDIGO
accent.line.fill.background()

add_textbox(s, Inches(1.2), Inches(2.0), Inches(11), Inches(0.6),
            "TEACHERTOOLS", font_size=16, bold=True, color=INDIGO)
add_textbox(s, Inches(1.2), Inches(2.6), Inches(11), Inches(1.5),
            "🗣️  AI 토론 수업 설계", font_size=54, bold=True, color=DARK)
add_textbox(s, Inches(1.2), Inches(4.1), Inches(11), Inches(1.2),
            "교사의 손을 덜고, 모든 학생에게 맞춤 토론 경험을",
            font_size=22, color=GRAY)
add_textbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4),
            "주제 한 줄 → 학습문제 · 지도안 · 토론 · 관찰 · 평가 · 생기부",
            font_size=14, color=INDIGO, bold=True)

# ── Slide 2 : 문제 제기 ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "토론 수업은 왜 어려운가", "교사가 가장 많이 꺾이는 지점")

# 3 problem cards
cards = [
    ("⏰", "준비 부담", "학습문제 · 지도안 · 평가지를 수업마다 새로 작성"),
    ("👀", "관찰 한계", "30명 학생 개별 발언을 실시간 기록·평가하기 어려움"),
    ("📝", "기록 부담", "학생별 생기부 문구 쓰다 보면 다음 수업 준비는 뒷전"),
]
card_w = Inches(3.9)
card_h = Inches(3.2)
gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_left = (Inches(13.333) - total_w) / 2
for i, (emoji, title, body) in enumerate(cards):
    left = start_left + (card_w + gap) * i
    top = Inches(2.0)
    add_rect(s, left, top, card_w, card_h, fill=WHITE, line_color=LIGHT_GRAY)
    add_textbox(s, left, top + Inches(0.3), card_w, Inches(0.8),
                emoji, font_size=44, align=PP_ALIGN.CENTER)
    add_textbox(s, left, top + Inches(1.4), card_w, Inches(0.5),
                title, font_size=22, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    add_textbox(s, left + Inches(0.3), top + Inches(2.0), card_w - Inches(0.6), Inches(1.1),
                body, font_size=15, color=GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.8),
            "➜   결국: “토론이 좋은 건 아는데, 현실에서 못 돌린다”",
            font_size=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

add_page_footer(s, 2, TOTAL)

# ── Slide 3 : 솔루션 한 줄 ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "솔루션", "토론 수업 전체 사이클을 AI가 자동화")

# 3-stage arrow
y = Inches(3.0)
box_w = Inches(3.2); box_h = Inches(1.6)
# box 1
b1 = add_rect(s, Inches(1.0), y, box_w, box_h, fill=INDIGO_LIGHT, line_color=INDIGO)
put_text_in_shape(b1, "수업 주제\n한 줄", font_size=22, bold=True, color=INDIGO_DARK)
# arrow
add_connector(s, Inches(4.3), y + box_h/2, Inches(5.1), y + box_h/2, color=INDIGO, width=3)
# box 2
b2 = add_rect(s, Inches(5.1), y, box_w, box_h, fill=INDIGO, line_color=INDIGO)
put_text_in_shape(b2, "AI 6단계\n오케스트레이션", font_size=22, bold=True, color=WHITE)
# arrow
add_connector(s, Inches(8.4), y + box_h/2, Inches(9.2), y + box_h/2, color=INDIGO, width=3)
# box 3
b3 = add_rect(s, Inches(9.2), y, box_w, box_h, fill=INDIGO_LIGHT, line_color=INDIGO)
put_text_in_shape(b3, "완성된 수업\n+ 학생별 생기부", font_size=22, bold=True, color=INDIGO_DARK)

add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.6),
            "교사는 “검토·승인”하는 자리로 이동. AI는 초안·관찰·평가를 준비.",
            font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.6),
            "최종 판단은 언제나 교사가 한다.",
            font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_page_footer(s, 3, TOTAL)

# ── Slide 4 : 6단계 플로우 ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "6단계 자동 생성", "각 단계마다 교사 승인 체크포인트")

steps = [
    ("①", "학습문제", "AI가 주제에 맞는\n3안 제안 → 선택"),
    ("②", "지도안·평가지", "지도안 · 루브릭 ·\n이원목적분류표"),
    ("③", "AI 토론", "학생 5명이 각자\nAI와 1:1 토론"),
    ("④", "과정 관찰", "학생별 5개 기준\n관찰 점수화"),
    ("⑤", "평가 판정", "상/중/하 + 총점\n+ 근거 인용"),
    ("⑥", "생기부 초안", "학생별 누가기록\n+ 교과학습발달상황"),
]
cols = 6
total_w = Inches(12.1)
step_w = Inches(1.85)
gap_w = (total_w - step_w * cols) / (cols - 1)
top = Inches(2.3)
for i, (num, title, body) in enumerate(steps):
    left = Inches(0.6) + (step_w + gap_w) * i
    fill = INDIGO if i in (2, 5) else WHITE  # highlight A3 토론, A6 생기부
    text_color = WHITE if fill == INDIGO else DARK
    num_color = WHITE if fill == INDIGO else INDIGO
    sub_color = WHITE if fill == INDIGO else GRAY
    add_rect(s, left, top, step_w, Inches(2.8), fill=fill, line_color=INDIGO, line_width=1.5)
    add_textbox(s, left, top + Inches(0.3), step_w, Inches(0.6),
                num, font_size=30, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_textbox(s, left + Inches(0.1), top + Inches(1.0), step_w - Inches(0.2), Inches(0.5),
                title, font_size=15, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    add_textbox(s, left + Inches(0.1), top + Inches(1.55), step_w - Inches(0.2), Inches(1.2),
                body, font_size=11, color=sub_color, align=PP_ALIGN.CENTER)

# approval check labels
add_textbox(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.5),
            "승인 게이트: ① 학습문제 선택,  ⑤ 평가 확인",
            font_size=16, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.6),
            "교사가 승인하지 않으면 다음 단계로 넘어가지 않음.",
            font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_page_footer(s, 4, TOTAL)

# ── Slide 5 : 핵심 기능 1 — 5명 동시 토론 ────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "학생 5명이 동시에, 각자 AI와", "핵심 기능 ①")

# left: bullets
add_bullets(s, Inches(0.8), Inches(2.1), Inches(6.2), Inches(4.5), [
    ("🔗", "교사는 QR 코드 5개를 한 번에 생성"),
    ("📱", "학생이 각자 폰으로 QR 스캔 → 본인 쓰레드 접속"),
    ("🎯", "학생은 찬성/반대 입장을 직접 선택하고 토론 시작"),
    ("👀", "교사는 한 화면에서 5명의 진행 상황 실시간 확인"),
    ("🔄", "학생은 “새로 시작하기”로 리셋 가능\n   (이전 대화는 교사 기록에 남음)"),
], font_size=17, line_spacing=1.5)

# right: mock QR dashboard
right_x = Inches(7.5)
right_y = Inches(2.1)
right_w = Inches(5.3)
right_h = Inches(4.3)
add_rect(s, right_x, right_y, right_w, right_h, fill=WHITE, line_color=LIGHT_GRAY)
add_textbox(s, right_x + Inches(0.3), right_y + Inches(0.3), right_w - Inches(0.6), Inches(0.4),
            "AI 토론 대시보드", font_size=13, bold=True, color=DARK)
# 5 QR cells in 2 rows (3+2)
cell_w = Inches(1.5); cell_h = Inches(1.45)
labels = [("#1 김민수", "찬성 · 진행중"), ("#2 이지은", "반대 · 완료"), ("#3 박서준", "찬성 · 진행중"),
          ("#4 최유나", "대기"), ("#5 정한결", "반대 · 진행중")]
pills = ["진행중", "완료", "진행중", "대기", "진행중"]
row_positions = [
    (right_x + Inches(0.3), right_y + Inches(0.8)),
    (right_x + Inches(1.9), right_y + Inches(0.8)),
    (right_x + Inches(3.5), right_y + Inches(0.8)),
    (right_x + Inches(0.3), right_y + Inches(2.4)),
    (right_x + Inches(1.9), right_y + Inches(2.4)),
]
for (lx, ly), (name, status), pill in zip(row_positions, labels, pills):
    add_rect(s, lx, ly, cell_w, cell_h, fill=BG, line_color=LIGHT_GRAY)
    # fake QR square
    qr_w = Inches(0.8); qr_h = Inches(0.8)
    qx = lx + (cell_w - qr_w) / 2
    qy = ly + Inches(0.2)
    add_rect(s, qx, qy, qr_w, qr_h, fill=DARK, line_color=DARK, rounded=False)
    add_textbox(s, lx, ly + Inches(1.05), cell_w, Inches(0.2),
                name, font_size=8, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # pill
    pill_color = {"진행중": INDIGO, "완료": GREEN, "대기": GRAY}[pill]
    add_textbox(s, lx, ly + Inches(1.22), cell_w, Inches(0.2),
                status, font_size=7, color=pill_color, align=PP_ALIGN.CENTER, bold=True)

add_page_footer(s, 5, TOTAL)

# ── Slide 6 : AI 오케스트레이터 ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "교사 옆에 AI 조교 한 명", "핵심 기능 ②: AI 오케스트레이터")

# left: bullets
add_bullets(s, Inches(0.8), Inches(2.1), Inches(6.5), Inches(4.6), [
    ("🧭", "5명 쓰레드 상태를 15초마다 스캔\n    → 요약 + 다음 행동 추천"),
    ("💬", "자연어 명령·피드백 채팅"),
    ("", "   “3,4번 토론 종료해줘”"),
    ("", "   “2번 태도 좋았어, 결과 상”"),
    ("🛠️", "AI가 해석 후 실제로 실행:"),
    ("", "   • 토론 종료 / 재시작"),
    ("", "   • 평가 등급 덮어쓰기"),
    ("", "   • 교사 메모 → 평가·생기부에 자동 반영"),
], font_size=16, line_spacing=1.35)

# right: chat mock
right_x = Inches(7.7)
right_y = Inches(2.1)
right_w = Inches(5.1)
right_h = Inches(4.3)
add_rect(s, right_x, right_y, right_w, right_h, fill=WHITE, line_color=INDIGO)

# header
add_textbox(s, right_x + Inches(0.3), right_y + Inches(0.2), Inches(4), Inches(0.4),
            "🧭 AI 오케스트레이터", font_size=13, bold=True, color=INDIGO_DARK)
add_textbox(s, right_x + Inches(0.3), right_y + Inches(0.55), right_w - Inches(0.6), Inches(0.4),
            "3명 완료, 2명 진행 중 (2번 3분째 입력 없음)",
            font_size=10, color=DARK)
add_textbox(s, right_x + Inches(0.3), right_y + Inches(0.9), right_w - Inches(0.6), Inches(0.3),
            "💡 2번 학생 상태 확인 / 나머지 완료 후 A4 실행",
            font_size=9, color=GRAY)

# chat bubbles
bub_y = right_y + Inches(1.5)
# teacher
t = add_rect(s, right_x + Inches(1.3), bub_y, Inches(3.5), Inches(0.45),
             fill=INDIGO, line_color=INDIGO)
put_text_in_shape(t, "3,4번 토론 종료해줘", font_size=11, color=WHITE, bold=False)
# AI
bub_y += Inches(0.6)
a = add_rect(s, right_x + Inches(0.3), bub_y, Inches(4.2), Inches(0.6),
             fill=WHITE, line_color=INDIGO_LIGHT)
put_text_in_shape(a, "3번, 4번 학생 토론을 종료했습니다.", font_size=10, color=DARK, bold=False, align=PP_ALIGN.LEFT)
# action pill
bub_y += Inches(0.7)
add_label_rect(s, right_x + Inches(0.3), bub_y, Inches(1.8), Inches(0.3),
               "토론 종료 · 3,4번", fill=INDIGO_LIGHT, text_color=INDIGO_DARK, font_size=9)
# teacher 2
bub_y += Inches(0.6)
t2 = add_rect(s, right_x + Inches(0.8), bub_y, Inches(4), Inches(0.45),
              fill=INDIGO, line_color=INDIGO)
put_text_in_shape(t2, "2번 태도 좋았어, 결과 상", font_size=11, color=WHITE)
# AI 2
bub_y += Inches(0.6)
a2 = add_rect(s, right_x + Inches(0.3), bub_y, Inches(4.2), Inches(0.45),
              fill=WHITE, line_color=INDIGO_LIGHT)
put_text_in_shape(a2, "2번 평가 상, 태도 메모 추가했습니다.", font_size=10, color=DARK)

add_page_footer(s, 6, TOTAL)

# ── Slide 7 : 교사 주도성 ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "AI는 도구, 최종 판단은 교사", "핵심 기능 ③: 교사 주도성")

points = [
    ("✏️", "루브릭 직접 편집", "AI가 만든 평가기준을 교사가 자신의 학급·수준에 맞게 수정 가능"),
    ("🔁", "재생성·다시 제안", "학습문제 마음에 안 들면 다시 제안, 평가도 재생성 가능"),
    ("🗣️", "정성 피드백 주입", "교사가 관찰한 정성적 평가를 AI 평가·생기부에 자동 반영"),
    ("✔", "승인 게이트", "학습문제 선택, 평가 판정 단계에서 교사 명시적 승인 필요"),
]
top = Inches(2.0)
for i, (emoji, title, body) in enumerate(points):
    y = top + Inches(1.15) * i
    # emoji circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.1), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = INDIGO_LIGHT
    circle.line.color.rgb = INDIGO_LIGHT
    put_text_in_shape(circle, emoji, font_size=22, bold=True, color=INDIGO_DARK)
    # title
    add_textbox(s, Inches(2.0), y + Inches(0.05), Inches(3.2), Inches(0.5),
                title, font_size=19, bold=True, color=DARK)
    # body
    add_textbox(s, Inches(5.4), y + Inches(0.15), Inches(7.4), Inches(0.9),
                body, font_size=15, color=GRAY)

add_page_footer(s, 7, TOTAL)

# ── Slide 8 : Agentic AI 아키텍처 ───────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "Agentic AI 아키텍처",
               "Perceive → Plan → Act → Observe 루프 기반 다중 에이전트 시스템")

# Coordinates
user_x, user_y, user_w, user_h = Inches(0.5), Inches(2.7), Inches(1.7), Inches(1.9)
agent_x, agent_y, agent_w, agent_h = Inches(2.55), Inches(1.7), Inches(5.0), Inches(4.3)
tools_x, tools_y, tools_w, tools_h = Inches(7.85), Inches(1.7), Inches(2.5), Inches(1.95)
sub_x, sub_y, sub_w, sub_h = Inches(7.85), Inches(3.85), Inches(2.5), Inches(2.15)
env_x, env_y, env_w, env_h = Inches(10.65), Inches(1.7), Inches(2.45), Inches(4.3)
mem_x, mem_y, mem_w, mem_h = Inches(0.5), Inches(6.2), Inches(12.6), Inches(0.75)

# USER block
add_rect(s, user_x, user_y, user_w, user_h, fill=INDIGO_LIGHT, line_color=INDIGO)
add_textbox(s, user_x, user_y + Inches(0.15), user_w, Inches(0.3),
            "USER", font_size=10, bold=True, color=INDIGO_DARK, align=PP_ALIGN.CENTER)
add_textbox(s, user_x, user_y + Inches(0.45), user_w, Inches(0.7),
            "👩‍🏫", font_size=40, align=PP_ALIGN.CENTER)
add_textbox(s, user_x, user_y + Inches(1.2), user_w, Inches(0.4),
            "교사", font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, user_x, user_y + Inches(1.55), user_w, Inches(0.3),
            "목표 · 감독", font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# AGENT core
add_rect(s, agent_x, agent_y, agent_w, agent_h, fill=WHITE, line_color=INDIGO, line_width=2.5)
add_textbox(s, agent_x, agent_y + Inches(0.18), agent_w, Inches(0.4),
            "🧭  ORCHESTRATOR AGENT (LLM)",
            font_size=14, bold=True, color=INDIGO_DARK, align=PP_ALIGN.CENTER)
add_textbox(s, agent_x, agent_y + Inches(0.6), agent_w, Inches(0.3),
            "자율적 의사결정 · 툴 호출 · 하위 에이전트 지휘",
            font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# 4-cycle nodes inside agent
nw, nh = Inches(2.1), Inches(1.15)
pad = Inches(0.25)
pl_x = agent_x + pad;                          pl_y = agent_y + Inches(1.05)
pr_x = agent_x + agent_w - pad - nw;           pr_y = pl_y
bl_x = pl_x;                                   bl_y = agent_y + agent_h - pad - nh
br_x = pr_x;                                   br_y = bl_y

pn = add_rect(s, pl_x, pl_y, nw, nh, fill=INDIGO_LIGHT, line_color=INDIGO)
put_text_in_shape(pn, "① PERCEIVE  지각\n상태·교사 메시지 읽기",
                  font_size=11, bold=True, color=INDIGO_DARK)
plan_n = add_rect(s, pr_x, pr_y, nw, nh, fill=INDIGO_LIGHT, line_color=INDIGO)
put_text_in_shape(plan_n, "② PLAN (LLM)  계획\n의도 분석 · 액션 선택",
                  font_size=11, bold=True, color=INDIGO_DARK)
act_n = add_rect(s, br_x, br_y, nw, nh, fill=INDIGO, line_color=INDIGO)
put_text_in_shape(act_n, "③ ACT  실행\n툴 호출 · 하위 에이전트",
                  font_size=11, bold=True, color=WHITE)
obs_n = add_rect(s, bl_x, bl_y, nw, nh, fill=INDIGO_LIGHT, line_color=INDIGO)
put_text_in_shape(obs_n, "④ OBSERVE  관찰\n결과 · 새 상태 확인",
                  font_size=11, bold=True, color=INDIGO_DARK)

# cycle arrows (clockwise)
add_connector(s, pl_x + nw, pl_y + nh/2, pr_x, pr_y + nh/2, color=INDIGO, width=2.0)
add_connector(s, pr_x + nw/2, pr_y + nh, br_x + nw/2, br_y, color=INDIGO, width=2.0)
add_connector(s, br_x, br_y + nh/2, bl_x + nw, bl_y + nh/2, color=INDIGO, width=2.0)
add_connector(s, bl_x + nw/2, bl_y, pl_x + nw/2, pl_y + nh, color=INDIGO, width=2.0)

# TOOLS box
add_rect(s, tools_x, tools_y, tools_w, tools_h, fill=WHITE, line_color=AMBER, line_width=1.5)
add_textbox(s, tools_x, tools_y + Inches(0.1), tools_w, Inches(0.35),
            "🛠️  TOOLS", font_size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_bullets(s, tools_x + Inches(0.25), tools_y + Inches(0.5), tools_w - Inches(0.35), Inches(1.4), [
    "finishThreads(idx)",
    "restartThreads(idx)",
    "addNote(idx, note)",
    "setGrade(idx, grade)",
], font_size=11, line_spacing=1.3)

# SUB-AGENTS box
add_rect(s, sub_x, sub_y, sub_w, sub_h, fill=WHITE, line_color=INDIGO, line_width=1.5)
add_textbox(s, sub_x, sub_y + Inches(0.1), sub_w, Inches(0.35),
            "🤖  SUB-AGENTS", font_size=13, bold=True, color=INDIGO_DARK, align=PP_ALIGN.CENTER)
add_bullets(s, sub_x + Inches(0.25), sub_y + Inches(0.5), sub_w - Inches(0.35), Inches(1.6), [
    "A1 학습문제",
    "A2 지도안 · 평가지",
    "A3 토론 × 5 (병렬)",
    "A4 관찰  /  A5 평가",
    "A6 생기부 초안",
], font_size=11, line_spacing=1.2)

# ENVIRONMENT box
add_rect(s, env_x, env_y, env_w, env_h, fill=BG, line_color=LIGHT_GRAY, line_width=1.5)
add_textbox(s, env_x, env_y + Inches(0.1), env_w, Inches(0.35),
            "🌐 ENVIRONMENT", font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, env_x, env_y + Inches(0.45), env_w, Inches(0.3),
            "학생 × 5 (모바일)", font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
for i in range(5):
    iy = env_y + Inches(0.85) + Inches(0.62) * i
    icon = add_rect(s, env_x + Inches(0.3), iy, env_w - Inches(0.6), Inches(0.5),
                    fill=WHITE, line_color=LIGHT_GRAY)
    put_text_in_shape(icon, f"📱  학생 {i+1}", font_size=11, color=DARK)

# MEMORY bar (bottom)
add_rect(s, mem_x, mem_y, mem_w, mem_h, fill=DARK, line_color=DARK)
add_textbox(s, mem_x + Inches(0.3), mem_y + Inches(0.1), Inches(3.5), Inches(0.55),
            "💾 MEMORY (shared)", font_size=13, bold=True, color=WHITE)
add_textbox(s, mem_x + Inches(3.9), mem_y + Inches(0.18), Inches(8.5), Inches(0.45),
            "세션 상태 · 쓰레드 대화 × 5 · 채팅 히스토리 · 교사 메모 · 평가 덮어쓰기 · 전체 대화 로그",
            font_size=11, color=LIGHT_GRAY)

# ── Arrows ──────────────────────────────────────────────────
# User ↔ Agent
add_connector(s, user_x + user_w, user_y + Inches(0.55),
              agent_x, agent_y + Inches(1.8), color=INDIGO, width=2.0, arrow=True)
add_connector(s, agent_x, agent_y + Inches(2.8),
              user_x + user_w, user_y + Inches(1.4), color=INDIGO, width=2.0, arrow=True)
add_textbox(s, user_x + user_w + Inches(0.02), agent_y + Inches(1.2),
            agent_x - user_x - user_w - Inches(0.04), Inches(0.3),
            "명령·피드백", font_size=9, color=INDIGO_DARK, align=PP_ALIGN.CENTER, bold=True)
add_textbox(s, user_x + user_w + Inches(0.02), agent_y + Inches(2.5),
            agent_x - user_x - user_w - Inches(0.04), Inches(0.3),
            "보고", font_size=9, color=INDIGO, align=PP_ALIGN.CENTER)

# Agent (ACT) → Tools
add_connector(s, agent_x + agent_w, agent_y + Inches(1.3), tools_x, tools_y + Inches(0.9),
              color=AMBER, width=2, arrow=True)
add_textbox(s, agent_x + agent_w + Inches(0.02), agent_y + Inches(1.0),
            tools_x - agent_x - agent_w - Inches(0.04), Inches(0.3),
            "invoke", font_size=9, color=AMBER, align=PP_ALIGN.CENTER, bold=True)

# Agent (ACT) → Sub-agents
add_connector(s, agent_x + agent_w, agent_y + Inches(3.2), sub_x, sub_y + Inches(1.0),
              color=INDIGO, width=2, arrow=True)
add_textbox(s, agent_x + agent_w + Inches(0.02), agent_y + Inches(3.0),
            sub_x - agent_x - agent_w - Inches(0.04), Inches(0.3),
            "delegate", font_size=9, color=INDIGO_DARK, align=PP_ALIGN.CENTER, bold=True)

# Sub-agents ↔ Environment (A3 converses with students)
add_connector(s, sub_x + sub_w, sub_y + Inches(1.2), env_x, env_y + Inches(2.1),
              color=GRAY, width=1.5, arrow=True)
add_connector(s, env_x, env_y + Inches(2.7), sub_x + sub_w, sub_y + Inches(1.6),
              color=GRAY, width=1.5, arrow=True)
add_textbox(s, sub_x + sub_w + Inches(0.01), sub_y + Inches(0.85),
            env_x - sub_x - sub_w - Inches(0.02), Inches(0.3),
            "A3 대화", font_size=9, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

# Agent ↔ Memory (read/write)
add_connector(s, pl_x + nw/2, agent_y + agent_h, pl_x + nw/2, mem_y,
              color=GREEN, width=1.5, arrow=True)
add_connector(s, bl_x + nw/2 + Inches(0.2), agent_y + agent_h, bl_x + nw/2 + Inches(0.2), mem_y,
              color=GREEN, width=1.5, arrow=True)
# Sub-agents → Memory
add_connector(s, sub_x + sub_w/2, sub_y + sub_h, sub_x + sub_w/2, mem_y,
              color=GREEN, width=1.5, arrow=True)

add_page_footer(s, 8, TOTAL)

# ── Slide 9 : 사용자 플로우 (교사 시간) ──────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "교사 시간 기준 한 차시 플로우", "실제 소요 시간 예시")

steps = [
    ("수업 전 (1분)", "주제 한 줄 입력 → 학습문제 3안 중 선택 후 승인"),
    ("수업 전 (2분)", "지도안·루브릭 검토, 필요 시 루브릭 항목·설명 수정"),
    ("수업 중 (15분)", "QR 5개 배포 → 학생 개별 토론\n오케스트레이터가 진행 상태·경고를 교사 화면에 표시"),
    ("수업 중 필요 시", "자연어 명령: “3번 종료”, “2번 결과 상” 등으로 실시간 개입"),
    ("수업 후 (2분)", "학생별 평가 탭 확인 → 필요 시 수정 → 생기부 초안 복사"),
]
top = Inches(2.0)
for i, (label, body) in enumerate(steps):
    y = top + Inches(0.95) * i
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = INDIGO
    circ.line.color.rgb = INDIGO
    put_text_in_shape(circ, str(i + 1), font_size=18, bold=True, color=WHITE)
    # label
    add_textbox(s, Inches(1.6), y, Inches(2.5), Inches(0.5),
                label, font_size=14, bold=True, color=INDIGO_DARK)
    # body
    add_textbox(s, Inches(4.2), y, Inches(8.7), Inches(0.9),
                body, font_size=15, color=DARK)

# highlight total
add_rect(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
         fill=AMBER, line_color=AMBER)
add_textbox(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.4),
            "교사가 직접 쓰는 시간: 학습문제 선택 + 승인 + 필요 시 조정  ≈  5분 내외",
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_page_footer(s, 9, TOTAL)

# ── Slide 10 : 기대 효과 ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "무엇이 달라지는가", "교사 경험 · 학생 경험 · 기록의 질")

impacts = [
    ("수업 준비 시간", "30분+", "5분 내외", "학습문제·지도안·평가지 자동화"),
    ("토론 빈도", "분기별 이벤트", "일상 수업", "준비 부담 제거"),
    ("평가 공정성", "교사 기억 의존", "AI 관찰 기록 + 교사 정성 판단", "근거 인용까지 자동 추출"),
    ("생기부 작성", "백지에서 시작", "학생별 초안에서 시작", "교사 메모도 자동 반영"),
]
col_x = [Inches(0.6), Inches(3.6), Inches(6.2), Inches(9.0)]
col_w = [Inches(2.9), Inches(2.5), Inches(2.7), Inches(3.7)]
headers = ["항목", "Before", "After", "핵심 변화"]
# header row
hy = Inches(2.0)
for x, w, h in zip(col_x, col_w, headers):
    add_rect(s, x, hy, w, Inches(0.55), fill=INDIGO, line_color=INDIGO)
    add_textbox(s, x, hy + Inches(0.1), w, Inches(0.4),
                h, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# rows
for i, row in enumerate(impacts):
    ry = hy + Inches(0.6) + Inches(0.95) * i
    bg = BG if i % 2 == 0 else WHITE
    for x, w, cell, is_after in zip(col_x, col_w, row, [False, False, True, False]):
        add_rect(s, x, ry, w, Inches(0.9), fill=bg, line_color=LIGHT_GRAY)
        color = INDIGO_DARK if is_after else DARK
        bold = True if is_after or col_x.index(x) == 0 else False
        add_textbox(s, x + Inches(0.15), ry + Inches(0.2), w - Inches(0.3), Inches(0.6),
                    cell, font_size=13, color=color, bold=bold, align=PP_ALIGN.CENTER)

add_page_footer(s, 10, TOTAL)

# ── Slide 11 : 로드맵 ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_header_bar(s, "로드맵", "MVP 이후 확장 계획")

phases = [
    ("MVP", "현재", "학생 5명 고정 / 단일 토론 / 기본 오케스트레이터", INDIGO),
    ("v1", "다음", "학생 수 가변 (2~30명) / 쓰레드 분할 배정", INDIGO),
    ("v1.5", "+α", "이전 시도 보기 · 학생별 승인 세분화", GRAY),
    ("v2", "중기", "학급 누적 데이터 / 학기 전체 생기부 합본", GRAY),
    ("v3", "장기", "토론 외 다양한 수업 템플릿 (프로젝트·발표·실험)", GRAY),
]
# timeline bar
top = Inches(2.3)
for i, (name, when, body, color) in enumerate(phases):
    y = top + Inches(0.85) * i
    # left pill
    add_rect(s, Inches(0.8), y, Inches(1.3), Inches(0.55), fill=color, line_color=color)
    add_textbox(s, Inches(0.8), y + Inches(0.1), Inches(1.3), Inches(0.4),
                name, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # when
    add_textbox(s, Inches(2.3), y + Inches(0.12), Inches(1.2), Inches(0.35),
                when, font_size=12, color=GRAY)
    # body
    add_textbox(s, Inches(3.7), y + Inches(0.1), Inches(9.2), Inches(0.5),
                body, font_size=15, color=DARK)

add_page_footer(s, 11, TOTAL)

# ── Slide 12 : End ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
# left bar
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.6), Inches(7.5))
accent.fill.solid(); accent.fill.fore_color.rgb = INDIGO
accent.line.fill.background()

add_textbox(s, Inches(1.2), Inches(2.3), Inches(11), Inches(1.2),
            "감사합니다.", font_size=60, bold=True, color=DARK)
add_textbox(s, Inches(1.2), Inches(3.7), Inches(11), Inches(0.6),
            "교사의 토론 수업을, 한 번 더 가능하게.",
            font_size=22, color=GRAY)
add_textbox(s, Inches(1.2), Inches(5.0), Inches(11), Inches(0.5),
            "Q & A",
            font_size=28, bold=True, color=INDIGO)
add_textbox(s, Inches(1.2), Inches(6.3), Inches(11), Inches(0.4),
            "TeacherTools · AI 토론 수업 설계",
            font_size=13, color=GRAY)

# ── Save ─────────────────────────────────────────────────────────
out = "docs/AI-토론-수업-설계.pptx"
prs.save(out)
print(f"saved: {out}")
print(f"slides: {len(prs.slides)}")
